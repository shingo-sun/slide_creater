"""
PowerPoint生成エンジン
解析結果からPPTXファイルを生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os
import logging
from typing import Dict, Any, Optional, List, Tuple
from pathlib import Path
from dataclasses import dataclass
from PIL import Image as PILImage

from app.models.slide_objects import (
    ObjectType, TextAlignment, Color
)

logger = logging.getLogger(__name__)


@dataclass
class TextBox:
    """テキストボックスの位置情報"""
    x: float
    y: float
    width: float
    height: float
    text: str
    font_size: float
    obj_data: Dict[str, Any]
    
    def intersects(self, other: 'TextBox', margin: float = 5) -> bool:
        """他のテキストボックスと重なっているか判定"""
        return not (
            self.x + self.width + margin < other.x or
            other.x + other.width + margin < self.x or
            self.y + self.height + margin < other.y or
            other.y + other.height + margin < self.y
        )
    
    def overlap_area(self, other: 'TextBox') -> float:
        """重なり面積を計算"""
        x_overlap = max(0, min(self.x + self.width, other.x + other.width) - max(self.x, other.x))
        y_overlap = max(0, min(self.y + self.height, other.y + other.height) - max(self.y, other.y))
        return x_overlap * y_overlap


class TextLayoutOptimizer:
    """テキスト配置最適化エンジン"""
    
    def __init__(self, slide_width: float, slide_height: float):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.placed_boxes: List[TextBox] = []
        self.shape_regions: List[Tuple[float, float, float, float]] = []
    
    def add_shape_region(self, x: float, y: float, width: float, height: float):
        """図形の占有領域を登録（テキストはこの領域外に配置優先）"""
        self.shape_regions.append((x, y, width, height))
    
    def optimize_text_objects(
        self, 
        text_objects: List[Dict[str, Any]],
        scale_x: float,
        scale_y: float
    ) -> List[Dict[str, Any]]:
        """
        テキストオブジェクトの配置を最適化
        - 重なり検出と位置調整
        - フォントサイズ調整
        """
        if not text_objects:
            return []
        
        # テキストを重要度・サイズでソート（大きいテキストを優先配置）
        sorted_texts = sorted(
            text_objects,
            key=lambda t: (
                -t.get("font_size", 12),  # 大きいフォント優先
                -len(t.get("text", "")),   # 長いテキスト優先
                t.get("y", 0)              # 上から順に
            )
        )
        
        optimized = []
        
        for obj in sorted_texts:
            optimized_obj = self._optimize_single_text(obj, scale_x, scale_y)
            if optimized_obj:
                optimized.append(optimized_obj)
        
        return optimized
    
    def _optimize_single_text(
        self,
        obj: Dict[str, Any],
        scale_x: float,
        scale_y: float
    ) -> Optional[Dict[str, Any]]:
        """単一テキストオブジェクトを最適化"""
        text = obj.get("text", "").strip()
        if not text:
            return None
        
        # 元の位置
        orig_x = obj["x"] * scale_x
        orig_y = obj["y"] * scale_y
        orig_width = obj["width"] * scale_x
        orig_height = obj["height"] * scale_y
        font_size = obj.get("font_size", 12)
        
        # テキストの長さに基づく最小サイズ計算
        min_width = self._estimate_text_width(text, font_size)
        min_height = self._estimate_text_height(text, font_size, orig_width)
        
        # サイズ調整
        new_width = max(orig_width, min_width)
        new_height = max(orig_height, min_height)
        
        # 新しいテキストボックス
        new_box = TextBox(
            x=orig_x,
            y=orig_y,
            width=new_width,
            height=new_height,
            text=text,
            font_size=font_size,
            obj_data=obj
        )
        
        # 重なり検出と位置調整
        adjusted_box = self._resolve_overlaps(new_box)
        
        # フォントサイズ調整（ボックスに収まらない場合）
        adjusted_font_size = self._adjust_font_size(
            text, adjusted_box.width, adjusted_box.height, font_size
        )
        
        # 結果を配置済みリストに追加
        self.placed_boxes.append(adjusted_box)
        
        # 最適化された座標を返す
        result = obj.copy()
        result["x"] = adjusted_box.x / scale_x
        result["y"] = adjusted_box.y / scale_y
        result["width"] = adjusted_box.width / scale_x
        result["height"] = adjusted_box.height / scale_y
        result["font_size"] = adjusted_font_size
        result["_optimized"] = True
        
        return result
    
    def _resolve_overlaps(self, box: TextBox, max_attempts: int = 10) -> TextBox:
        """重なりを解消する位置を探す"""
        if not self.placed_boxes:
            return box
        
        current_box = TextBox(
            x=box.x, y=box.y, width=box.width, height=box.height,
            text=box.text, font_size=box.font_size, obj_data=box.obj_data
        )
        
        for attempt in range(max_attempts):
            overlapping = [b for b in self.placed_boxes if current_box.intersects(b)]
            
            if not overlapping:
                return current_box
            
            # 最も重なりが大きいボックスを見つける
            worst_overlap = max(overlapping, key=lambda b: current_box.overlap_area(b))
            
            # 移動方向を決定（右、下、左、上の順で試す）
            directions = [
                (worst_overlap.x + worst_overlap.width + 10, current_box.y),  # 右
                (current_box.x, worst_overlap.y + worst_overlap.height + 10),  # 下
                (worst_overlap.x - current_box.width - 10, current_box.y),     # 左
                (current_box.x, worst_overlap.y - current_box.height - 10),    # 上
            ]
            
            best_pos = None
            min_distance = float('inf')
            
            for new_x, new_y in directions:
                # 境界チェック
                if new_x < 0 or new_y < 0:
                    continue
                if new_x + current_box.width > self.slide_width:
                    continue
                if new_y + current_box.height > self.slide_height:
                    continue
                
                # 元の位置からの距離
                distance = ((new_x - box.x) ** 2 + (new_y - box.y) ** 2) ** 0.5
                
                if distance < min_distance:
                    min_distance = distance
                    best_pos = (new_x, new_y)
            
            if best_pos:
                current_box.x, current_box.y = best_pos
            else:
                # 縮小して対応
                current_box.width *= 0.9
                current_box.height *= 0.9
        
        return current_box
    
    def _estimate_text_width(self, text: str, font_size: float) -> float:
        """テキスト幅を推定（EMU単位）"""
        # 日本語: 1文字 ≈ フォントサイズ
        # 英数字: 1文字 ≈ フォントサイズ * 0.6
        jp_chars = sum(1 for c in text if ord(c) > 127)
        en_chars = len(text) - jp_chars
        
        char_width = Pt(font_size).emu
        estimated = (jp_chars * char_width) + (en_chars * char_width * 0.6)
        
        return estimated
    
    def _estimate_text_height(self, text: str, font_size: float, max_width: float) -> float:
        """テキスト高さを推定（折り返し考慮）"""
        text_width = self._estimate_text_width(text, font_size)
        lines = max(1, int(text_width / max_width) + 1)
        line_height = Pt(font_size).emu * 1.5
        return lines * line_height
    
    def _adjust_font_size(
        self, 
        text: str, 
        box_width: float, 
        box_height: float, 
        original_size: float,
        min_size: float = 8
    ) -> float:
        """ボックスに収まるようフォントサイズを調整"""
        current_size = original_size
        
        while current_size > min_size:
            text_width = self._estimate_text_width(text, current_size)
            text_height = self._estimate_text_height(text, current_size, box_width)
            
            if text_width <= box_width * 1.1 and text_height <= box_height * 1.1:
                return current_size
            
            current_size -= 1
        
        return min_size


class PPTXGenerator:
    """PowerPoint生成エンジン"""
    
    # スライドサイズ定義（16:9固定、EMU単位）
    # 標準の16:9スライド: 13.333インチ x 7.5インチ
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    ASPECT_RATIO = 16 / 9
    
    # スライドサイズ定義（互換性のため残す）
    SLIDE_SIZES = {
        "16:9": (Inches(13.333), Inches(7.5)),
        "4:3": (Inches(10), Inches(7.5)),
    }
    
    def __init__(self):
        self.prs = None
        self.slide = None
        self.image_width = 0
        self.image_height = 0
        self.scale_x = 1.0
        self.scale_y = 1.0
        # 16:9固定
        self.slide_width = self.SLIDE_WIDTH
        self.slide_height = self.SLIDE_HEIGHT
    
    def generate(
        self, 
        analysis_result: Dict[str, Any],
        output_path: str,
        original_image_path: Optional[str] = None,
        use_background_image: bool = False,
        hybrid_mode: bool = True  # 新パラメータ: ハイブリッドモード
    ) -> str:
        """
        解析結果からPPTXを生成
        
        Args:
            analysis_result: 解析結果辞書
            output_path: 出力ファイルパス
            original_image_path: 元画像パス（フォールバック用）
            use_background_image: 元画像を背景として使用するか
            hybrid_mode: Trueの場合、図形は画像として配置、テキストのみ編集可能
            
        Returns:
            生成されたPPTXのパス
        """
        # プレゼンテーション作成
        self.prs = Presentation()
        
        # スライドサイズ設定（16:9固定）
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        
        # 画像サイズを取得
        image_info = analysis_result.get("image_info", {})
        self.image_width = image_info.get("width", 1920)
        self.image_height = image_info.get("height", 1080)
        
        # スケール計算（アスペクト比を維持してフィット）
        self._calculate_scale_with_aspect_ratio()
        
        # ブランクスライド追加
        blank_layout = self.prs.slide_layouts[6]  # ブランクレイアウト
        self.slide = self.prs.slides.add_slide(blank_layout)
        
        # モード判定
        mode = analysis_result.get("mode", "hybrid" if hybrid_mode else "faithful")
        
        if mode == "hybrid" or hybrid_mode:
            # ハイブリッドモード: 図形は画像、テキストは編集可能
            self._generate_hybrid(analysis_result, original_image_path)
        elif use_background_image and original_image_path and os.path.exists(original_image_path):
            # 背景画像モード
            self._add_background_image(original_image_path, self.slide_width, self.slide_height)
            self._add_editable_text_overlay(analysis_result)
        else:
            # 通常モード: すべてを再構築
            self._generate_standard(analysis_result)
        
        # 保存
        self.prs.save(output_path)
        logger.info(f"PPTX生成完了: {output_path}")
        
        return output_path
    
    def _calculate_scale_with_aspect_ratio(self):
        """アスペクト比を維持してスケールを計算（縮尺エラーを防ぐ）"""
        image_aspect = self.image_width / self.image_height
        slide_aspect = self.ASPECT_RATIO
        
        if abs(image_aspect - slide_aspect) < 0.01:
            # ほぼ同じアスペクト比の場合
            self.scale_x = self.slide_width / self.image_width
            self.scale_y = self.slide_height / self.image_height
        elif image_aspect > slide_aspect:
            # 画像の方が横長 → 幅に合わせる
            self.scale_x = self.slide_width / self.image_width
            self.scale_y = self.scale_x  # 同じスケールで縦横比維持
        else:
            # 画像の方が縦長 → 高さに合わせる
            self.scale_y = self.slide_height / self.image_height
            self.scale_x = self.scale_y  # 同じスケールで縦横比維持
        
        logger.debug(f"スケール計算: x={self.scale_x}, y={self.scale_y}")
    
    def _generate_hybrid(self, analysis_result: Dict[str, Any], original_image_path: Optional[str]):
        """
        ハイブリッドモードで生成（最も確実な方式）
        - 元画像をそのまま背景として配置（見た目100%一致）
        - テキストのみ編集可能なテキストボックスとしてオーバーレイ
        """
        # 1. 元画像を背景として配置（スライド全体にフィット）
        if original_image_path and os.path.exists(original_image_path):
            self._add_full_background_image(original_image_path)
            logger.info(f"背景画像を配置: {original_image_path}")
        else:
            logger.warning("元画像が見つかりません。図形のみ配置します。")
            # フォールバック: 切り抜き画像を配置
            cropped_images = analysis_result.get("cropped_images", [])
            for img_obj in cropped_images:
                self._add_cropped_image(img_obj)
        
        # 2. テキストを編集可能なテキストボックスとして配置（最前面）
        text_objects = analysis_result.get("text_objects", [])
        
        # テキストをソート（大きいフォント・上から順に）
        sorted_texts = sorted(
            [t for t in text_objects if t.get("text", "").strip()],
            key=lambda t: (-t.get("font_size", 12), t.get("y", 0))
        )
        
        placed_count = 0
        for obj in sorted_texts:
            text = obj.get("text", "").strip()
            if not text:
                continue
            
            # テキストボックスを追加
            self._add_text_overlay(obj)
            placed_count += 1
        
        logger.info(
            f"ハイブリッドモード生成完了: "
            f"背景画像1枚, テキスト{placed_count}個"
        )
    
    def _add_full_background_image(self, image_path: str):
        """元画像をスライド全体の背景として配置"""
        # スライド全体にフィットするように配置（最もシンプルで確実）
        try:
            self.slide.shapes.add_picture(
                image_path,
                0, 0,  # 左上から
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT  # スライド全体
            )
            logger.debug(f"背景画像を配置: {image_path}")
        except Exception as e:
            logger.error(f"背景画像配置エラー: {e}")
    
    def _add_text_overlay(self, obj: Dict[str, Any]):
        """編集可能な透明テキストボックスをオーバーレイ"""
        text = obj.get("text", "").strip()
        if not text:
            return
        
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        # 最小サイズ確保
        min_width = Pt(30).emu
        min_height = Pt(15).emu
        width = max(width, min_width)
        height = max(height, min_height)
        
        try:
            # テキストボックス追加
            textbox = self.slide.shapes.add_textbox(
                Emu(x), Emu(y), Emu(width), Emu(height)
            )
            
            # 背景を透明に
            textbox.fill.background()
            
            # テキストフレーム設定
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(1).emu
            tf.margin_right = Pt(1).emu
            tf.margin_top = Pt(0).emu
            tf.margin_bottom = Pt(0).emu
            
            # 段落設定
            p = tf.paragraphs[0]
            p.text = text
            
            # フォント設定
            font = p.font
            font_size = obj.get("font_size", 12)
            font_size = max(6, min(72, font_size))
            font.size = Pt(font_size)
            font.name = obj.get("font_family", "Meiryo")
            
            # 太字
            if obj.get("font_weight") == "bold":
                font.bold = True
            
            # 文字色（透明にするか、元のままか選択可能）
            text_color = obj.get("text_color")
            if text_color:
                # 完全透明にして背景画像のテキストを見せる
                # または元の色を維持
                font.color.rgb = RGBColor(
                    text_color.get("r", 0),
                    text_color.get("g", 0),
                    text_color.get("b", 0)
                )
            
            # アライメント
            alignment = obj.get("alignment", "left")
            if alignment == "center":
                p.alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
        except Exception as e:
            logger.error(f"テキストオーバーレイ配置エラー: {e}")
    
    def _generate_standard(self, analysis_result: Dict[str, Any]):
        """通常モードで生成（従来のロジック）"""
        optimizer = TextLayoutOptimizer(self.slide_width, self.slide_height)
        
        # 図形オブジェクトの領域を登録
        shape_objects = analysis_result.get("shape_objects", [])
        for obj in shape_objects:
            x, y, w, h = self._scale_position(
                obj["x"], obj["y"], obj["width"], obj["height"]
            )
            optimizer.add_shape_region(x, y, w, h)
        
        # テキストを最適化
        text_objects = analysis_result.get("text_objects", [])
        optimized_texts = optimizer.optimize_text_objects(
            text_objects, self.scale_x, self.scale_y
        )
        
        optimized_result = analysis_result.copy()
        optimized_result["text_objects"] = optimized_texts
        
        all_objects = self._collect_and_sort_objects(optimized_result)
        
        for obj in all_objects:
            self._add_object(obj, "faithful")
    
    def _add_cropped_image(self, obj: Dict[str, Any]):
        """切り抜き画像をスライドに配置"""
        image_path = obj.get("image_path")
        if not image_path or not os.path.exists(image_path):
            logger.warning(f"画像ファイルが見つかりません: {image_path}")
            return
        
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        try:
            self.slide.shapes.add_picture(
                image_path,
                Emu(x), Emu(y), Emu(width), Emu(height)
            )
            logger.debug(f"切り抜き画像を配置: {image_path} at ({x}, {y})")
        except Exception as e:
            logger.error(f"画像配置エラー: {e}")
    
    def _add_background_image(
        self, 
        image_path: str, 
        slide_width: int, 
        slide_height: int
    ):
        """元画像を背景として追加"""
        # スライド全体に画像を配置
        self.slide.shapes.add_picture(
            image_path,
            Emu(0),
            Emu(0),
            slide_width,
            slide_height
        )
    
    def _add_editable_text_overlay(self, analysis_result: Dict[str, Any]):
        """
        背景画像の上に編集可能なテキストボックスをオーバーレイ
        （透明背景で元画像のテキスト位置に配置）
        """
        text_objects = analysis_result.get("text_objects", [])
        
        # 信頼度の高いテキストのみ、最大10個まで
        sorted_texts = sorted(
            [t for t in text_objects if t.get("text", "").strip()],
            key=lambda t: t.get("source_confidence", 0),
            reverse=True
        )[:10]
        
        for obj in sorted_texts:
            if not obj.get("text", "").strip():
                continue
                
            x, y, width, height = self._scale_position(
                obj["x"], obj["y"], obj["width"], obj["height"]
            )
            
            # 透明なテキストボックスを追加
            textbox = self.slide.shapes.add_textbox(
                Emu(x), Emu(y), Emu(width), Emu(height)
            )
            
            # 背景を透明に
            textbox.fill.background()
            
            # テキスト設定
            tf = textbox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = obj.get("text", "")
            
            font = p.font
            font.size = Pt(obj.get("font_size", 12))
            font.name = obj.get("font_family", "Meiryo")
            
            # 透明テキスト（ユーザーが編集時に見える）
            # font.color.rgb = RGBColor(0, 0, 0)  # コメントアウト：背景と重ねるため
    
    def _collect_and_sort_objects(
        self, analysis_result: Dict[str, Any]
    ) -> List[Dict[str, Any]]:
        """全オブジェクトを収集してz_index順にソート"""
        all_objects = []
        
        # 図形オブジェクト
        for obj in analysis_result.get("shape_objects", []):
            all_objects.append(obj)
        
        # アイコンオブジェクト
        for obj in analysis_result.get("icon_objects", []):
            all_objects.append(obj)
        
        # テキストオブジェクト（最前面）
        for obj in analysis_result.get("text_objects", []):
            all_objects.append(obj)
        
        # z_indexでソート
        all_objects.sort(key=lambda x: x.get("z_index", 0))
        
        return all_objects
    
    def _add_object(self, obj: Dict[str, Any], mode: str):
        """オブジェクトをスライドに追加"""
        obj_type = obj.get("object_type")
        
        if obj_type == "text":
            self._add_text_box(obj)
        elif obj_type == "rectangle":
            self._add_rectangle(obj)
        elif obj_type == "rounded_rectangle":
            self._add_rounded_rectangle(obj)
        elif obj_type == "circle":
            self._add_circle(obj)
        elif obj_type == "ellipse":
            self._add_ellipse(obj)
        elif obj_type == "line":
            self._add_line(obj)
        elif obj_type == "arrow":
            self._add_arrow(obj)
        elif obj_type == "icon":
            self._add_icon(obj)
        elif obj_type == "image":
            self._add_image(obj)
    
    def _scale_position(
        self, x: float, y: float, width: float, height: float
    ) -> tuple:
        """座標をスライドサイズにスケーリング"""
        return (
            int(x * self.scale_x),
            int(y * self.scale_y),
            int(width * self.scale_x),
            int(height * self.scale_y)
        )
    
    def _add_text_box(self, obj: Dict[str, Any]):
        """テキストボックスを追加（最適化対応）"""
        text = obj.get("text", "").strip()
        if not text:
            return
        
        # 最適化済みの場合はスケーリング不要
        if obj.get("_optimized"):
            x = int(obj["x"] * self.scale_x)
            y = int(obj["y"] * self.scale_y)
            width = int(obj["width"] * self.scale_x)
            height = int(obj["height"] * self.scale_y)
        else:
            x, y, width, height = self._scale_position(
                obj["x"], obj["y"], obj["width"], obj["height"]
            )
        
        # 最小サイズを確保
        min_width = Pt(50).emu
        min_height = Pt(20).emu
        width = max(width, min_width)
        height = max(height, min_height)
        
        # テキストボックス追加
        textbox = self.slide.shapes.add_textbox(
            Emu(x), Emu(y), Emu(width), Emu(height)
        )
        
        # テキストフレーム設定
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # 内部マージンを小さく
        tf.margin_left = Pt(2).emu
        tf.margin_right = Pt(2).emu
        tf.margin_top = Pt(1).emu
        tf.margin_bottom = Pt(1).emu
        
        # 段落設定
        p = tf.paragraphs[0]
        p.text = text
        
        # フォント設定
        font = p.font
        font_size = obj.get("font_size", 12)
        
        # フォントサイズの妥当性チェック
        if font_size > 72:
            font_size = 72
        elif font_size < 6:
            font_size = 6
        
        font.size = Pt(font_size)
        font.name = obj.get("font_family", "Meiryo")
        
        # 太字
        if obj.get("font_weight") == "bold":
            font.bold = True
        
        # 文字色
        text_color = obj.get("text_color")
        if text_color:
            font.color.rgb = RGBColor(
                text_color.get("r", 0),
                text_color.get("g", 0),
                text_color.get("b", 0)
            )
        
        # アライメント
        alignment = obj.get("alignment", "left")
        if alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
            
        logger.debug(f"テキスト配置: '{text[:20]}...' at ({x}, {y}) size {font_size}pt")
    
    def _add_rectangle(self, obj: Dict[str, Any]):
        """長方形を追加"""
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(x), Emu(y), Emu(width), Emu(height)
        )
        
        self._apply_shape_style(shape, obj)
    
    def _add_rounded_rectangle(self, obj: Dict[str, Any]):
        """角丸長方形を追加"""
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x), Emu(y), Emu(width), Emu(height)
        )
        
        # 角丸の調整
        corner_radius = obj.get("corner_radius", 10)
        # python-pptxでは adjustments で角丸を設定
        if shape.adjustments:
            try:
                shape.adjustments[0] = min(0.5, corner_radius / min(width, height))
            except:
                pass
        
        self._apply_shape_style(shape, obj)
    
    def _add_circle(self, obj: Dict[str, Any]):
        """円を追加"""
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(x), Emu(y), Emu(width), Emu(height)
        )
        
        self._apply_shape_style(shape, obj)
    
    def _add_ellipse(self, obj: Dict[str, Any]):
        """楕円を追加"""
        self._add_circle(obj)  # 同じ処理
    
    def _add_line(self, obj: Dict[str, Any]):
        """線を追加"""
        start_x = int(obj.get("start_x", obj["x"]) * self.scale_x)
        start_y = int(obj.get("start_y", obj["y"]) * self.scale_y)
        end_x = int(obj.get("end_x", obj["x"] + obj["width"]) * self.scale_x)
        end_y = int(obj.get("end_y", obj["y"] + obj["height"]) * self.scale_y)
        
        # コネクタとして追加
        connector = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(start_x), Emu(start_y),
            Emu(end_x), Emu(end_y)
        )
        
        # 線のスタイル
        line = connector.line
        line_color = obj.get("line_color")
        if line_color:
            line.color.rgb = RGBColor(
                line_color.get("r", 0),
                line_color.get("g", 0),
                line_color.get("b", 0)
            )
        
        line.width = Pt(obj.get("line_width", 1))
    
    def _add_arrow(self, obj: Dict[str, Any]):
        """矢印を追加"""
        start_x = int(obj.get("start_x", obj["x"]) * self.scale_x)
        start_y = int(obj.get("start_y", obj["y"]) * self.scale_y)
        end_x = int(obj.get("end_x", obj["x"] + obj["width"]) * self.scale_x)
        end_y = int(obj.get("end_y", obj["y"] + obj["height"]) * self.scale_y)
        
        # 矢印付きコネクタ
        connector = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(start_x), Emu(start_y),
            Emu(end_x), Emu(end_y)
        )
        
        # 矢印スタイル設定
        line = connector.line
        line_color = obj.get("line_color")
        if line_color:
            line.color.rgb = RGBColor(
                line_color.get("r", 0),
                line_color.get("g", 0),
                line_color.get("b", 0)
            )
        
        line.width = Pt(obj.get("line_width", 2))
        
        # 矢印ヘッドを設定（XMLレベルで設定）
        arrow_head = obj.get("arrow_head", "end")
        self._set_arrow_head(connector, arrow_head)
    
    def _set_arrow_head(self, connector, arrow_head: str):
        """矢印の先端を設定"""
        try:
            line = connector.line._ln
            
            if arrow_head in ["end", "both"]:
                # 終端に矢印
                tailEnd = parse_xml(
                    '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
                )
                line.append(tailEnd)
            
            if arrow_head in ["start", "both"]:
                # 始端に矢印
                headEnd = parse_xml(
                    '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
                )
                line.append(headEnd)
        except Exception as e:
            print(f"矢印設定エラー: {e}")
    
    def _add_icon(self, obj: Dict[str, Any]):
        """アイコン画像を追加"""
        image_path = obj.get("image_path")
        if not image_path or not os.path.exists(image_path):
            return
        
        x, y, width, height = self._scale_position(
            obj["x"], obj["y"], obj["width"], obj["height"]
        )
        
        try:
            self.slide.shapes.add_picture(
                image_path,
                Emu(x), Emu(y), Emu(width), Emu(height)
            )
        except Exception as e:
            print(f"アイコン追加エラー: {e}")
    
    def _add_image(self, obj: Dict[str, Any]):
        """画像を追加（フォールバック用）"""
        self._add_icon(obj)
    
    def _apply_shape_style(self, shape, obj: Dict[str, Any]):
        """図形にスタイルを適用"""
        # 塗りつぶし色
        fill_color = obj.get("fill_color")
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(
                fill_color.get("r", 255),
                fill_color.get("g", 255),
                fill_color.get("b", 255)
            )
        else:
            shape.fill.background()  # 塗りつぶしなし
        
        # 線色
        line_color = obj.get("line_color")
        if line_color:
            shape.line.color.rgb = RGBColor(
                line_color.get("r", 0),
                line_color.get("g", 0),
                line_color.get("b", 0)
            )
            shape.line.width = Pt(obj.get("line_width", 1))
        else:
            shape.line.fill.background()  # 線なし
        
        # 透明度
        transparency = obj.get("transparency", 0)
        if transparency > 0:
            try:
                # 透明度の設定（0-100%）
                fill = shape.fill._xPr
                if fill is not None:
                    solidFill = fill.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is not None:
                            alpha = parse_xml(
                                f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{int((1-transparency/100)*100000)}"/>'
                            )
                            srgbClr.append(alpha)
            except:
                pass
    
    def add_background_image(
        self, 
        image_path: str, 
        opacity: float = 0.3
    ):
        """背景画像を追加（忠実再現モードのガイド用）"""
        if not os.path.exists(image_path):
            return
        
        # スライドサイズで配置
        picture = self.slide.shapes.add_picture(
            image_path,
            Emu(0), Emu(0),
            self.prs.slide_width, self.prs.slide_height
        )
        
        # 最背面に移動
        spTree = self.slide.shapes._spTree
        sp = picture._element
        spTree.remove(sp)
        spTree.insert(2, sp)  # 背景の次に配置
