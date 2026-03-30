"""
PPTX→画像変換サービス
生成したPPTXを画像に変換して比較可能にする
"""
import os
import subprocess
import tempfile
import shutil
from pathlib import Path
from typing import Optional, List, Tuple
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import logging

logger = logging.getLogger(__name__)


class PPTXToImageConverter:
    """PPTXを画像に変換するサービス"""
    
    def __init__(self):
        self.libreoffice_path = self._find_libreoffice()
        
    def _find_libreoffice(self) -> Optional[str]:
        """LibreOfficeのパスを検索"""
        # macOS
        mac_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/usr/local/bin/soffice",
        ]
        # Linux
        linux_paths = [
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
        ]
        
        for path in mac_paths + linux_paths:
            if os.path.exists(path):
                return path
        
        # PATH から検索
        try:
            result = subprocess.run(
                ["which", "soffice"],
                capture_output=True,
                text=True
            )
            if result.returncode == 0:
                return result.stdout.strip()
        except Exception:
            pass
            
        return None
    
    def convert_pptx_to_images(
        self,
        pptx_path: str,
        output_dir: str,
        format: str = "jpg",
        dpi: int = 150
    ) -> List[str]:
        """
        PPTXを画像に変換
        
        Args:
            pptx_path: PPTXファイルパス
            output_dir: 出力ディレクトリ
            format: 出力フォーマット (jpg, png)
            dpi: 解像度
            
        Returns:
            生成された画像パスのリスト
        """
        os.makedirs(output_dir, exist_ok=True)
        
        # LibreOfficeがあれば使用
        if self.libreoffice_path:
            try:
                return self._convert_with_libreoffice(pptx_path, output_dir, format)
            except Exception as e:
                logger.warning(f"LibreOffice変換失敗、フォールバック使用: {e}")
        
        # フォールバック: python-pptx + PILでレンダリング
        logger.info("python-pptx + PILでPPTXをレンダリングします")
        return self._convert_with_pil(pptx_path, output_dir, format, dpi)
    
    def _convert_with_pil(
        self,
        pptx_path: str,
        output_dir: str,
        format: str,
        dpi: int
    ) -> List[str]:
        """python-pptx + PILでPPTXを画像にレンダリング"""
        try:
            prs = Presentation(pptx_path)
            output_paths = []
            
            # スライドサイズ取得
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # ピクセルに変換（EMUからピクセル）
            # 1 inch = 914400 EMU, dpi pixels per inch
            scale = dpi / 914400.0
            img_width = int(slide_width * scale)
            img_height = int(slide_height * scale)
            
            # 最小サイズ確保
            img_width = max(img_width, 960)
            img_height = max(img_height, 540)
            
            for i, slide in enumerate(prs.slides):
                # 白背景の画像を作成
                img = Image.new('RGB', (img_width, img_height), 'white')
                draw = ImageDraw.Draw(img)
                
                # スケール係数
                x_scale = img_width / slide_width
                y_scale = img_height / slide_height
                
                # 各シェイプをレンダリング
                for shape in slide.shapes:
                    self._render_shape(draw, shape, x_scale, y_scale, img_width, img_height)
                
                # 保存
                output_path = os.path.join(output_dir, f"slide_{i + 1}.{format}")
                if format.lower() == "jpg":
                    img = img.convert("RGB")
                    img.save(output_path, "JPEG", quality=95)
                else:
                    img.save(output_path, format.upper())
                
                output_paths.append(output_path)
                logger.info(f"スライド {i + 1} をレンダリング: {output_path}")
            
            return output_paths
            
        except Exception as e:
            logger.error(f"PILレンダリングエラー: {e}")
            raise Exception(f"PPTXレンダリング失敗: {e}")
    
    def _render_shape(
        self,
        draw: ImageDraw.Draw,
        shape,
        x_scale: float,
        y_scale: float,
        img_width: int,
        img_height: int
    ):
        """シェイプをレンダリング"""
        try:
            # 位置とサイズを取得
            left = int(shape.left * x_scale) if shape.left else 0
            top = int(shape.top * y_scale) if shape.top else 0
            width = int(shape.width * x_scale) if shape.width else 100
            height = int(shape.height * y_scale) if shape.height else 100
            
            # 境界チェック
            left = max(0, min(left, img_width - 10))
            top = max(0, min(top, img_height - 10))
            right = min(left + width, img_width)
            bottom = min(top + height, img_height)
            
            # テキストフレームがある場合
            if shape.has_text_frame:
                self._render_text(draw, shape.text_frame, left, top, right, bottom)
            
            # 図形の場合
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                self._render_auto_shape(draw, shape, left, top, right, bottom)
            
            # グループの場合
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    self._render_shape(draw, subshape, x_scale, y_scale, img_width, img_height)
            
            # 画像の場合
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._render_picture(draw, shape, left, top, right, bottom)
            
            # テーブルの場合
            elif shape.has_table:
                self._render_table(draw, shape.table, left, top, right, bottom)
                
        except Exception as e:
            logger.debug(f"シェイプレンダリングスキップ: {e}")
    
    def _render_text(
        self,
        draw: ImageDraw.Draw,
        text_frame,
        left: int,
        top: int,
        right: int,
        bottom: int
    ):
        """テキストをレンダリング"""
        try:
            y_pos = top + 5
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    
                    # フォントサイズ取得
                    font_size = 14
                    if run.font.size:
                        font_size = max(8, min(48, int(run.font.size.pt)))
                    
                    # 色取得
                    text_color = (0, 0, 0)  # デフォルト黒
                    try:
                        if run.font.color and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            text_color = (rgb[0], rgb[1], rgb[2])
                    except:
                        pass
                    
                    # フォント
                    try:
                        font = ImageFont.truetype("/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc", font_size)
                    except:
                        try:
                            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
                        except:
                            font = ImageFont.load_default()
                    
                    # テキスト描画
                    draw.text((left + 5, y_pos), text, fill=text_color, font=font)
                    y_pos += font_size + 5
                    
        except Exception as e:
            logger.debug(f"テキストレンダリングエラー: {e}")
    
    def _render_auto_shape(
        self,
        draw: ImageDraw.Draw,
        shape,
        left: int,
        top: int,
        right: int,
        bottom: int
    ):
        """オートシェイプをレンダリング"""
        try:
            # 塗りつぶし色取得
            fill_color = (200, 200, 200)  # デフォルトグレー
            try:
                if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    rgb = shape.fill.fore_color.rgb
                    fill_color = (rgb[0], rgb[1], rgb[2])
            except:
                pass
            
            # 線の色取得
            outline_color = (100, 100, 100)
            try:
                if shape.line and shape.line.color and shape.line.color.rgb:
                    rgb = shape.line.color.rgb
                    outline_color = (rgb[0], rgb[1], rgb[2])
            except:
                pass
            
            # 矩形を描画
            draw.rectangle([left, top, right, bottom], fill=fill_color, outline=outline_color)
            
            # テキストがあれば描画
            if shape.has_text_frame:
                self._render_text(draw, shape.text_frame, left, top, right, bottom)
                
        except Exception as e:
            logger.debug(f"オートシェイプレンダリングエラー: {e}")
    
    def _render_picture(
        self,
        draw: ImageDraw.Draw,
        shape,
        left: int,
        top: int,
        right: int,
        bottom: int
    ):
        """画像プレースホルダーを描画"""
        # 画像領域を示すプレースホルダー
        draw.rectangle([left, top, right, bottom], fill=(230, 230, 230), outline=(180, 180, 180))
        draw.line([(left, top), (right, bottom)], fill=(180, 180, 180), width=1)
        draw.line([(left, bottom), (right, top)], fill=(180, 180, 180), width=1)
    
    def _render_table(
        self,
        draw: ImageDraw.Draw,
        table,
        left: int,
        top: int,
        right: int,
        bottom: int
    ):
        """テーブルをレンダリング"""
        try:
            rows = len(table.rows)
            cols = len(table.columns)
            
            if rows == 0 or cols == 0:
                return
            
            cell_width = (right - left) // cols
            cell_height = (bottom - top) // rows
            
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    # セル位置
                    cell_left = left + col_idx * cell_width
                    cell_top = top + row_idx * cell_height
                    cell_right = cell_left + cell_width
                    cell_bottom = cell_top + cell_height
                    
                    # セル枠描画
                    draw.rectangle(
                        [cell_left, cell_top, cell_right, cell_bottom],
                        fill=(255, 255, 255),
                        outline=(0, 0, 0)
                    )
                    
                    # セルテキスト
                    if cell.text_frame:
                        self._render_text(draw, cell.text_frame, cell_left, cell_top, cell_right, cell_bottom)
                        
        except Exception as e:
            logger.debug(f"テーブルレンダリングエラー: {e}")
    
    def _convert_with_libreoffice(
        self,
        pptx_path: str,
        output_dir: str,
        format: str
    ) -> List[str]:
        """LibreOfficeで変換"""
        try:
            # まずPDFに変換
            with tempfile.TemporaryDirectory() as temp_dir:
                cmd = [
                    self.libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", temp_dir,
                    pptx_path
                ]
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                if result.returncode != 0:
                    logger.error(f"LibreOffice変換エラー: {result.stderr}")
                    raise Exception(f"LibreOffice変換失敗: {result.stderr}")
                
                # PDFファイルを探す
                pdf_files = list(Path(temp_dir).glob("*.pdf"))
                if not pdf_files:
                    raise Exception("PDFファイルが生成されませんでした")
                
                pdf_path = str(pdf_files[0])
                
                # PDFを画像に変換
                return self._pdf_to_images(pdf_path, output_dir, format)
                
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice変換がタイムアウトしました")
            raise Exception("変換がタイムアウトしました")
    
    def _pdf_to_images(
        self,
        pdf_path: str,
        output_dir: str,
        format: str,
        dpi: int = 150
    ) -> List[str]:
        """PDFを画像に変換"""
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(pdf_path, dpi=dpi)
            output_paths = []
            
            for i, image in enumerate(images):
                output_path = os.path.join(
                    output_dir,
                    f"slide_{i + 1}.{format}"
                )
                
                if format.lower() == "jpg":
                    image = image.convert("RGB")
                    image.save(output_path, "JPEG", quality=95)
                else:
                    image.save(output_path, format.upper())
                
                output_paths.append(output_path)
            
            return output_paths
            
        except ImportError:
            logger.warning("pdf2imageがインストールされていません")
            raise Exception("pdf2imageが必要です: pip install pdf2image")
    
    def convert_single_slide(
        self,
        pptx_path: str,
        slide_index: int = 0,
        format: str = "jpg"
    ) -> str:
        """
        単一スライドを画像に変換
        
        Args:
            pptx_path: PPTXファイルパス
            slide_index: スライド番号（0始まり）
            format: 出力フォーマット
            
        Returns:
            画像ファイルパス
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            images = self.convert_pptx_to_images(
                pptx_path, temp_dir, format
            )
            
            if slide_index >= len(images):
                raise ValueError(f"スライド {slide_index} は存在しません")
            
            # 永続的な場所にコピー
            output_dir = os.path.dirname(pptx_path)
            output_path = os.path.join(
                output_dir,
                f"generated_slide_{slide_index + 1}.{format}"
            )
            
            shutil.copy(images[slide_index], output_path)
            
            return output_path


# シングルトンインスタンス
pptx_to_image_converter = PPTXToImageConverter()
