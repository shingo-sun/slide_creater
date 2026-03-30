"""
画像忠実再現型 PowerPoint再構成AIシステム - バックエンドAPI
"""
# 環境変数を.envから読み込み
from dotenv import load_dotenv
load_dotenv()

from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
import os
import uuid
import shutil
import cv2
import logging
from pathlib import Path
from typing import Optional
from enum import Enum

from app.services.image_analyzer import ImageAnalyzer
from app.services.ocr_service import OCRService
from app.services.shape_detector import ShapeDetector
from app.services.pptx_generator import PPTXGenerator
from app.services.layout_analyzer import LayoutAnalyzer
from app.services.quality_checker import QualityChecker
from app.services.shape_cropper import shape_cropper
from app.models.slide_objects import SlideAnalysisResult, ConversionMode

logger = logging.getLogger(__name__)

app = FastAPI(
    title="Slide Creater API",
    description="画像からPowerPointを忠実再現するAIシステム",
    version="1.0.0"
)

# CORS設定
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ディレクトリ設定（絶対パス）
BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# 静的ファイル配信
app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")
app.mount("/outputs", StaticFiles(directory=str(OUTPUT_DIR)), name="outputs")


class ConversionModeEnum(str, Enum):
    faithful = "faithful"  # 忠実再現モード
    editable = "editable"  # 編集優先モード
    hybrid = "hybrid"      # ハイブリッドモード


@app.get("/")
async def root():
    return {"message": "Slide Creater API", "version": "1.0.0"}


@app.post("/api/upload")
async def upload_image(file: UploadFile = File(...)):
    """画像ファイルをアップロード"""
    # ファイル形式チェック
    allowed_extensions = {".png", ".jpg", ".jpeg", ".pdf"}
    file_ext = Path(file.filename).suffix.lower()
    
    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"サポートされていないファイル形式です: {file_ext}"
        )
    
    # ユニークなファイル名を生成
    file_id = str(uuid.uuid4())
    filename = f"{file_id}{file_ext}"
    file_path = UPLOAD_DIR / filename
    
    # ファイル保存
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    return {
        "file_id": file_id,
        "filename": filename,
        "original_name": file.filename,
        "path": str(file_path)
    }


@app.post("/api/analyze/{file_id}")
async def analyze_image(
    file_id: str,
    mode: ConversionModeEnum = ConversionModeEnum.faithful,
    slide_ratio: str = "16:9"
):
    """画像を解析してオブジェクト情報を抽出（品質チェック付き）"""
    # ファイル検索
    file_path = None
    for ext in [".png", ".jpg", ".jpeg", ".pdf"]:
        candidate = UPLOAD_DIR / f"{file_id}{ext}"
        if candidate.exists():
            file_path = candidate
            break
    
    if not file_path:
        raise HTTPException(status_code=404, detail="ファイルが見つかりません")
    
    try:
        # 画像解析パイプライン
        image_analyzer = ImageAnalyzer()
        ocr_service = OCRService()
        shape_detector = ShapeDetector()
        layout_analyzer = LayoutAnalyzer()
        quality_checker = QualityChecker()
        
        # 1. 前処理
        processed_image, image_info = image_analyzer.preprocess(str(file_path))
        original_image = cv2.imread(str(file_path))
        
        # 2. レイアウト解析
        layout_blocks = layout_analyzer.analyze(processed_image)
        
        # 3. OCR処理
        text_objects = ocr_service.extract_text(str(file_path), processed_image)
        
        # 4. 図形検出
        shape_objects = shape_detector.detect(processed_image)
        
        # 5. アイコン検出・切り出し
        icon_objects = shape_detector.detect_icons(processed_image, UPLOAD_DIR, file_id)
        
        # 初回結果統合
        result = SlideAnalysisResult(
            file_id=file_id,
            image_info=image_info,
            layout_blocks=layout_blocks,
            text_objects=text_objects,
            shape_objects=shape_objects,
            icon_objects=icon_objects,
            mode=mode.value,
            slide_ratio=slide_ratio
        )
        
        # 6. 品質チェック（ダブルループ検証）
        quality_report = quality_checker.validate_analysis_result(
            original_image, result
        )
        
        # 7. 品質が低い場合は最適化を適用
        if not quality_report.is_acceptable:
            print(f"品質チェック: スコア {quality_report.overall_score:.2f}")
            print(f"問題点: {quality_report.issues}")
            print(f"推奨: {quality_report.recommendations}")
            
            # 解析結果を最適化
            result = quality_checker.optimize_analysis_result(result, quality_report)
        
        # 品質情報を追加
        result_dict = result.model_dump()
        result_dict["quality_report"] = {
            "is_acceptable": quality_report.is_acceptable,
            "overall_score": quality_report.overall_score,
            "text_coverage_score": quality_report.text_coverage_score,
            "layout_score": quality_report.layout_score,
            "issues": quality_report.issues,
            "recommendations": quality_report.recommendations
        }
        
        return result_dict
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"解析エラー: {str(e)}")


@app.post("/api/generate/{file_id}")
async def generate_pptx(
    file_id: str,
    mode: ConversionModeEnum = ConversionModeEnum.faithful,
    slide_ratio: str = "16:9",
    use_background_image: Optional[bool] = None,
    hybrid_mode: bool = True  # 新パラメータ: ハイブリッドモード
):
    """
    解析結果からPPTXを生成
    
    hybrid_mode=True (デフォルト):
    - 元画像をそのまま背景として配置
    - テキストは編集可能なテキストボックスとしてオーバーレイ
    - 縮尺エラーを防ぐため16:9固定
    """
    # まず解析を実行
    analysis_response = await analyze_image(file_id, mode, slide_ratio)
    
    try:
        # 元画像パス取得
        original_image_path = None
        for ext in [".png", ".jpg", ".jpeg"]:
            candidate = UPLOAD_DIR / f"{file_id}{ext}"
            if candidate.exists():
                original_image_path = str(candidate)
                break
        
        if not original_image_path:
            raise HTTPException(status_code=404, detail="元画像が見つかりません")
        
        # 品質レポートを確認
        quality_report = analysis_response.get("quality_report", {})
        
        # PPTX生成
        pptx_generator = PPTXGenerator()
        output_filename = f"{file_id}.pptx"
        output_path = OUTPUT_DIR / output_filename
        
        pptx_generator.generate(
            analysis_result=analysis_response,
            output_path=str(output_path),
            original_image_path=original_image_path,
            use_background_image=False,  # ハイブリッドモードで処理
            hybrid_mode=hybrid_mode
        )
        
        return {
            "file_id": file_id,
            "pptx_path": str(output_path),
            "download_url": f"/api/download/{file_id}",
            "hybrid_mode": hybrid_mode,
            "text_objects_count": len(analysis_response.get("text_objects", [])),
            "quality_report": quality_report,
            "analysis": analysis_response
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPTX生成エラー: {str(e)}")


@app.get("/api/download/{file_id}")
async def download_pptx(file_id: str):
    """生成されたPPTXをダウンロード"""
    output_path = OUTPUT_DIR / f"{file_id}.pptx"
    
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="PPTXファイルが見つかりません")
    
    return FileResponse(
        path=str(output_path),
        filename=f"slide_{file_id}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@app.get("/api/compare/{file_id}")
async def compare_images(file_id: str):
    """元画像と生成結果の比較情報を取得"""
    # 元画像パス
    original_path = None
    for ext in [".png", ".jpg", ".jpeg"]:
        candidate = UPLOAD_DIR / f"{file_id}{ext}"
        if candidate.exists():
            original_path = f"/uploads/{file_id}{ext}"
            break
    
    if not original_path:
        raise HTTPException(status_code=404, detail="元画像が見つかりません")
    
    # PPTX存在確認
    pptx_path = OUTPUT_DIR / f"{file_id}.pptx"
    pptx_exists = pptx_path.exists()
    
    return {
        "file_id": file_id,
        "original_image_url": original_path,
        "pptx_exists": pptx_exists,
        "pptx_download_url": f"/api/download/{file_id}" if pptx_exists else None
    }


@app.delete("/api/cleanup/{file_id}")
async def cleanup_files(file_id: str):
    """一時ファイルを削除"""
    deleted_files = []
    
    # アップロードファイル削除
    for ext in [".png", ".jpg", ".jpeg", ".pdf"]:
        file_path = UPLOAD_DIR / f"{file_id}{ext}"
        if file_path.exists():
            file_path.unlink()
            deleted_files.append(str(file_path))
    
    # アイコン画像削除
    for icon_file in UPLOAD_DIR.glob(f"{file_id}_icon_*"):
        icon_file.unlink()
        deleted_files.append(str(icon_file))
    
    # PPTX削除
    pptx_path = OUTPUT_DIR / f"{file_id}.pptx"
    if pptx_path.exists():
        pptx_path.unlink()
        deleted_files.append(str(pptx_path))
    
    return {"deleted_files": deleted_files}


# ============================================================
# AI分析・学習・自動修正 エンドポイント
# ============================================================

from app.services.pptx_to_image import pptx_to_image_converter
from app.services.image_comparator import image_comparator
from app.services.ai_analyzer import ai_image_analyzer
from app.services.learning_store import learning_data_store
from app.services.auto_editor import auto_editor
from app.services.layer_decomposer import layer_decomposer, DecompositionResult


@app.post("/api/validate/{file_id}")
async def validate_generation(file_id: str):
    """
    生成結果を検証: PPTX→画像変換→比較→AI分析
    ダブルループ学習の核となるエンドポイント
    """
    # 元画像パス取得
    original_path = None
    for ext in [".png", ".jpg", ".jpeg"]:
        candidate = UPLOAD_DIR / f"{file_id}{ext}"
        if candidate.exists():
            original_path = str(candidate)
            break
    
    if not original_path:
        raise HTTPException(status_code=404, detail="元画像が見つかりません")
    
    # PPTX確認
    pptx_path = OUTPUT_DIR / f"{file_id}.pptx"
    if not pptx_path.exists():
        raise HTTPException(status_code=404, detail="PPTXが生成されていません")
    
    try:
        # 1. PPTX → 画像変換
        generated_image_dir = OUTPUT_DIR / f"{file_id}_images"
        generated_images = pptx_to_image_converter.convert_pptx_to_images(
            str(pptx_path),
            str(generated_image_dir),
            format="jpg"
        )
        
        if not generated_images:
            raise HTTPException(status_code=500, detail="PPTX画像変換に失敗しました")
        
        generated_path = generated_images[0]
        
        # 2. 画像比較
        diff_image_path = str(OUTPUT_DIR / f"{file_id}_diff.jpg")
        comparison_result = image_comparator.compare_images(
            original_path,
            generated_path,
            diff_image_path
        )
        
        # 3. AI分析（Claude API）
        ai_result = await ai_image_analyzer.analyze_difference(
            original_path,
            generated_path,
            comparison_result.summary
        )
        
        # 4. 学習データに保存
        error_pattern_ids = []
        for diff in ai_result.differences:
            pattern = learning_data_store.add_error_pattern(
                pattern_type=diff.get("type", "unknown"),
                description=diff.get("original", ""),
                original_state={"description": diff.get("original", "")},
                incorrect_state={"description": diff.get("generated", "")},
                correct_state={"description": diff.get("original", "")},
                prevention_rule=None
            )
            error_pattern_ids.append(pattern.id)
        
        # 学習エントリ追加
        learning_data_store.add_learning_entry(
            original_image_path=original_path,
            similarity_score=comparison_result.similarity_score,
            error_pattern_ids=error_pattern_ids,
            ai_feedback={
                "quality": ai_result.quality_assessment,
                "feedback": ai_result.detailed_feedback,
                "insights": ai_result.learning_insights
            },
            corrections_applied=[],
            result_improvement=0.0
        )
        
        return {
            "file_id": file_id,
            "validation": {
                "similarity_score": comparison_result.similarity_score,
                "ssim_score": comparison_result.ssim_score,
                "pixel_diff_ratio": comparison_result.pixel_diff_ratio,
                "difference_regions": [
                    {
                        "x": r.x, "y": r.y,
                        "width": r.width, "height": r.height,
                        "type": r.type, "severity": r.severity,
                        "description": r.description
                    }
                    for r in comparison_result.difference_regions
                ],
                "summary": comparison_result.summary
            },
            "ai_analysis": {
                "quality_assessment": ai_result.quality_assessment,
                "differences": ai_result.differences,
                "corrections": ai_result.corrections,
                "feedback": ai_result.detailed_feedback,
                "learning_insights": ai_result.learning_insights,
                "confidence": ai_result.confidence
            },
            "generated_image_url": f"/outputs/{file_id}_images/slide_1.jpg",
            "diff_image_url": f"/outputs/{file_id}_diff.jpg"
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"検証エラー: {str(e)}")


@app.post("/api/regenerate/{file_id}")
async def regenerate_with_corrections(
    file_id: str,
    apply_ai_corrections: bool = True
):
    """
    AI分析結果を反映して再生成
    学習データを活用した改善版PPTX生成
    """
    # 元画像パス取得
    original_path = None
    for ext in [".png", ".jpg", ".jpeg"]:
        candidate = UPLOAD_DIR / f"{file_id}{ext}"
        if candidate.exists():
            original_path = str(candidate)
            break
    
    if not original_path:
        raise HTTPException(status_code=404, detail="元画像が見つかりません")
    
    try:
        # 1. 解析実行
        analysis_result = await analyze_image(file_id)
        
        # 2. 学習データからルールを適用
        analysis_result = learning_data_store.apply_learned_rules(analysis_result)
        applied_rules = analysis_result.pop("_applied_learning_rules", [])
        
        # 3. AI修正提案を適用（オプション）
        applied_corrections = []
        if apply_ai_corrections:
            # 前回の検証結果からAI提案を取得
            stats = learning_data_store.get_statistics()
            patterns = learning_data_store.get_relevant_patterns(
                min_frequency=1,
                min_confidence=0.5
            )
            
            # パターンから修正提案を生成
            ai_corrections = []
            for pattern in patterns[:5]:  # 上位5パターン
                if pattern.prevention_rule:
                    ai_corrections.append({
                        "target": pattern.pattern_type,
                        "action": "apply_rule",
                        "parameters": {"rule": pattern.prevention_rule},
                        "priority": 2,
                        "description": pattern.description
                    })
            
            if ai_corrections:
                analysis_result = auto_editor.apply_corrections(
                    analysis_result,
                    ai_corrections
                )
                applied_corrections = analysis_result.pop("_applied_corrections", [])
        
        # 4. PPTX再生成
        pptx_generator = PPTXGenerator()
        output_filename = f"{file_id}_v2.pptx"
        output_path = OUTPUT_DIR / output_filename
        
        pptx_generator.generate(
            analysis_result=analysis_result,
            output_path=str(output_path),
            original_image_path=original_path,
            use_background_image=False
        )
        
        return {
            "file_id": file_id,
            "version": 2,
            "pptx_path": str(output_path),
            "download_url": f"/outputs/{output_filename}",
            "applied_learning_rules": applied_rules,
            "applied_corrections": applied_corrections,
            "message": "学習データを活用して再生成しました"
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"再生成エラー: {str(e)}")


@app.get("/api/learning/statistics")
async def get_learning_statistics():
    """学習データの統計情報を取得"""
    return learning_data_store.get_statistics()


@app.get("/api/learning/patterns")
async def get_error_patterns(
    pattern_type: Optional[str] = None,
    min_frequency: int = 1
):
    """登録されているエラーパターンを取得"""
    patterns = learning_data_store.get_relevant_patterns(
        pattern_type=pattern_type,
        min_frequency=min_frequency
    )
    
    return {
        "total": len(patterns),
        "patterns": [
            {
                "id": p.id,
                "type": p.pattern_type,
                "description": p.description,
                "frequency": p.frequency,
                "confidence": p.confidence,
                "prevention_rule": p.prevention_rule,
                "last_occurred": p.last_occurred
            }
            for p in patterns
        ]
    }


@app.post("/api/learning/add-rule")
async def add_prevention_rule(
    pattern_id: str,
    rule: str
):
    """エラーパターンに防止ルールを追加"""
    if pattern_id not in learning_data_store.error_patterns:
        raise HTTPException(status_code=404, detail="パターンが見つかりません")
    
    pattern = learning_data_store.error_patterns[pattern_id]
    pattern.prevention_rule = rule
    pattern.confidence = min(pattern.confidence + 0.1, 1.0)
    learning_data_store._save_data()
    
    return {
        "pattern_id": pattern_id,
        "rule": rule,
        "message": "防止ルールを追加しました"
    }


@app.post("/api/full-pipeline/{file_id}")
async def full_pipeline_with_validation(
    file_id: str,
    max_iterations: int = 5,
    target_similarity: float = 0.95
):
    """
    完全自動パイプライン：
    生成 → 検証 → 学習 → 再生成 を目標品質(95%)に達するまで繰り返す
    
    ダブルループ学習のフルサイクル:
    Step 1: 初回PPTX生成
    Step 2: 画像比較検証 (PPTX→画像変換→比較)
    Step 3: Claude API分析 (差分の詳細分析)
    Step 4: 学習データ蓄積
    Step 5: AI修正適用・再生成
    → Step 2-5を95%達成まで繰り返し
    """
    iterations_data = []
    current_similarity = 0.0
    best_similarity = 0.0
    best_version = 1
    
    for iteration in range(max_iterations):
        iteration_result = {
            "iteration": iteration + 1,
            "steps": [],
            "similarity_score": 0.0,
            "status": "in_progress"
        }
        
        try:
            # ========== Step 1: 生成 ==========
            iteration_result["steps"].append({
                "step": 1,
                "name": "PPTX生成",
                "status": "running"
            })
            
            if iteration == 0:
                gen_result = await generate_pptx(file_id)
                iteration_result["steps"][-1]["description"] = "初回PPTX生成"
            else:
                gen_result = await regenerate_with_corrections(file_id)
                iteration_result["steps"][-1]["description"] = f"AI修正を適用して再生成 (v{iteration + 1})"
            
            iteration_result["steps"][-1]["status"] = "completed"
            iteration_result["pptx_url"] = gen_result.get("download_url")
            
            # ========== Step 2: PPTX→画像変換 ==========
            iteration_result["steps"].append({
                "step": 2,
                "name": "PPTX画像化",
                "status": "running",
                "description": "生成したPPTXを画像に変換"
            })
            
            # 元画像パス取得
            original_path = None
            for ext in [".png", ".jpg", ".jpeg"]:
                candidate = UPLOAD_DIR / f"{file_id}{ext}"
                if candidate.exists():
                    original_path = str(candidate)
                    break
            
            pptx_path = OUTPUT_DIR / f"{file_id}.pptx"
            generated_image_dir = OUTPUT_DIR / f"{file_id}_images"
            generated_images = pptx_to_image_converter.convert_pptx_to_images(
                str(pptx_path),
                str(generated_image_dir),
                format="jpg"
            )
            
            if generated_images:
                generated_path = generated_images[0]
                iteration_result["steps"][-1]["status"] = "completed"
                iteration_result["generated_image_url"] = f"/outputs/{file_id}_images/slide_1.jpg"
            else:
                iteration_result["steps"][-1]["status"] = "failed"
                iteration_result["steps"][-1]["error"] = "画像変換失敗"
                iteration_result["status"] = "failed"
                iterations_data.append(iteration_result)
                continue
            
            # ========== Step 3: 画像比較検証 ==========
            iteration_result["steps"].append({
                "step": 3,
                "name": "画像比較",
                "status": "running",
                "description": "元画像と生成画像のピクセル・構造比較"
            })
            
            diff_image_path = str(OUTPUT_DIR / f"{file_id}_diff_v{iteration + 1}.jpg")
            comparison_result = image_comparator.compare_images(
                original_path,
                generated_path,
                diff_image_path
            )
            
            current_similarity = comparison_result.similarity_score
            iteration_result["similarity_score"] = current_similarity
            iteration_result["ssim_score"] = comparison_result.ssim_score
            iteration_result["steps"][-1]["status"] = "completed"
            iteration_result["steps"][-1]["result"] = {
                "similarity": f"{current_similarity * 100:.1f}%",
                "ssim": f"{comparison_result.ssim_score * 100:.1f}%"
            }
            iteration_result["diff_image_url"] = f"/outputs/{file_id}_diff_v{iteration + 1}.jpg"
            
            # ベストスコア更新
            if current_similarity > best_similarity:
                best_similarity = current_similarity
                best_version = iteration + 1
            
            # ========== 目標達成チェック ==========
            if current_similarity >= target_similarity:
                iteration_result["steps"].append({
                    "step": 4,
                    "name": "目標達成",
                    "status": "completed",
                    "description": f"類似度 {current_similarity * 100:.1f}% ≥ 目標 {target_similarity * 100:.0f}%"
                })
                iteration_result["status"] = "target_achieved"
                iterations_data.append(iteration_result)
                
                return {
                    "status": "success",
                    "message": f"🎉 目標品質 {target_similarity * 100:.0f}% を達成しました！",
                    "final_similarity": current_similarity,
                    "iterations_count": iteration + 1,
                    "best_version": best_version,
                    "iterations": iterations_data,
                    "download_url": gen_result.get("download_url")
                }
            
            # ========== Step 4: Claude API分析 ==========
            iteration_result["steps"].append({
                "step": 4,
                "name": "Claude AI分析",
                "status": "running",
                "description": "Claude APIで差分を詳細分析・修正提案生成"
            })
            
            ai_result = await ai_image_analyzer.analyze_difference(
                original_path,
                generated_path,
                comparison_result.summary
            )
            
            iteration_result["steps"][-1]["status"] = "completed"
            iteration_result["ai_analysis"] = {
                "quality_assessment": ai_result.quality_assessment,
                "differences_count": len(ai_result.differences),
                "corrections_count": len(ai_result.corrections),
                "confidence": ai_result.confidence
            }
            iteration_result["differences"] = ai_result.differences
            iteration_result["corrections"] = ai_result.corrections
            
            # ========== Step 5: 学習データ蓄積 ==========
            iteration_result["steps"].append({
                "step": 5,
                "name": "学習データ保存",
                "status": "running",
                "description": "エラーパターンと修正方法を学習データに蓄積"
            })
            
            for diff in ai_result.differences:
                learning_data_store.add_error_pattern(
                    pattern_type=diff.get("type", "unknown"),
                    description=diff.get("original", ""),
                    original_state={"description": diff.get("original", "")},
                    incorrect_state={"description": diff.get("generated", "")},
                    correct_state={"description": diff.get("original", "")},
                    prevention_rule=None
                )
            
            iteration_result["steps"][-1]["status"] = "completed"
            iteration_result["status"] = "completed"
            
            # 修正提案がない場合は終了
            if not ai_result.corrections:
                iteration_result["steps"].append({
                    "step": 6,
                    "name": "修正完了",
                    "status": "completed",
                    "description": "これ以上の修正提案がありません"
                })
                iterations_data.append(iteration_result)
                break
                
        except Exception as e:
            import traceback
            traceback.print_exc()
            iteration_result["status"] = "error"
            iteration_result["error"] = str(e)
            if iteration_result["steps"]:
                iteration_result["steps"][-1]["status"] = "failed"
                iteration_result["steps"][-1]["error"] = str(e)
        
        iterations_data.append(iteration_result)
    
    # 最大イテレーション到達または終了
    return {
        "status": "incomplete" if best_similarity < target_similarity else "success",
        "message": f"最終類似度: {best_similarity * 100:.1f}%（目標: {target_similarity * 100:.0f}%）",
        "final_similarity": best_similarity,
        "iterations_count": len(iterations_data),
        "best_version": best_version,
        "iterations": iterations_data,
        "download_url": f"/api/download/{file_id}"
    }


# ============================================================
# Magic Layers風 画像分解エンドポイント
# ============================================================

@app.post("/api/decompose/{file_id}")
async def decompose_image(file_id: str):
    """
    Magic Layers風画像分解:
    画像を編集可能なレイヤーに分解し、各要素を透過PNGとして出力
    
    処理フロー:
    1. Claude Vision APIで画像内の要素を検出
    2. OpenCVで境界を精密化
    3. 各要素を透過PNGとして切り出し
    4. メタデータJSON出力
    
    Returns:
        分解されたレイヤー一覧、各PNG URL、メタデータ
    """
    # ファイル検索
    file_path = None
    for ext in [".png", ".jpg", ".jpeg"]:
        candidate = UPLOAD_DIR / f"{file_id}{ext}"
        if candidate.exists():
            file_path = str(candidate)
            break
    
    if not file_path:
        raise HTTPException(status_code=404, detail="画像ファイルが見つかりません")
    
    try:
        # レイヤー分解実行
        result = await layer_decomposer.decompose(file_path, file_id)
        
        if not result.success:
            raise HTTPException(status_code=500, detail=result.message)
        
        # レイヤー情報をレスポンス用に整形
        layers_response = []
        for layer in result.layers:
            # 画像URLを生成
            image_url = None
            if layer.image_path:
                relative_path = Path(layer.image_path).relative_to(BASE_DIR)
                image_url = f"/{relative_path}"
            
            layers_response.append({
                "layer_id": layer.layer_id,
                "element_type": layer.element_type,
                "x": layer.x,
                "y": layer.y,
                "width": layer.width,
                "height": layer.height,
                "z_index": layer.z_index,
                "description": layer.description,
                "content": layer.content,
                "color": layer.color,
                "image_url": image_url,
                "confidence": layer.confidence
            })
        
        return {
            "file_id": file_id,
            "success": True,
            "message": result.message,
            "original_size": {
                "width": result.original_width,
                "height": result.original_height
            },
            "background_color": result.background_color,
            "layer_count": len(result.layers),
            "layers": layers_response,
            "output_dir": f"/outputs/{file_id}_layers",
            "metadata_url": f"/outputs/{file_id}_layers/layers_metadata.json"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"レイヤー分解エラー: {str(e)}")


@app.get("/api/layers/{file_id}")
async def get_decomposed_layers(file_id: str):
    """
    分解済みレイヤー情報を取得
    """
    metadata_path = OUTPUT_DIR / f"{file_id}_layers" / "layers_metadata.json"
    
    if not metadata_path.exists():
        raise HTTPException(
            status_code=404, 
            detail="レイヤーデータが見つかりません。先に /api/decompose/{file_id} を実行してください。"
        )
    
    import json
    with open(metadata_path, "r", encoding="utf-8") as f:
        metadata = json.load(f)
    
    # 画像URLを追加
    for layer in metadata.get("layers", []):
        if layer.get("image_path"):
            relative_path = Path(layer["image_path"]).relative_to(BASE_DIR)
            layer["image_url"] = f"/{relative_path}"
    
    return {
        "file_id": file_id,
        **metadata
    }


@app.delete("/api/layers/{file_id}")
async def delete_decomposed_layers(file_id: str):
    """
    分解済みレイヤーを削除
    """
    layers_dir = OUTPUT_DIR / f"{file_id}_layers"
    
    if not layers_dir.exists():
        raise HTTPException(status_code=404, detail="レイヤーデータが見つかりません")
    
    import shutil
    deleted_count = len(list(layers_dir.glob("*.png")))
    shutil.rmtree(layers_dir)
    
    return {
        "file_id": file_id,
        "deleted_layers": deleted_count,
        "message": "レイヤーデータを削除しました"
    }


@app.post("/api/layers/{file_id}/pptx")
async def generate_pptx_from_layers(file_id: str):
    """
    分解されたレイヤーからPPTXを生成
    - テキスト要素 → 編集可能なテキストボックス
    - 図形（四角、線など）→ PowerPointネイティブシェイプ
    - 画像/アイコン → PNG画像として配置
    """
    import json
    import tempfile
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    import io
    
    # メタデータ読み込み
    metadata_path = OUTPUT_DIR / f"{file_id}_layers" / "layers_metadata.json"
    
    if not metadata_path.exists():
        raise HTTPException(
            status_code=404, 
            detail="レイヤーデータが見つかりません。先に /api/decompose/{file_id} を実行してください。"
        )
    
    with open(metadata_path, "r", encoding="utf-8") as f:
        metadata = json.load(f)
    
    def hex_to_rgb(hex_color: str) -> RGBColor:
        """HEXカラーをRGBColorに変換"""
        if not hex_color or not hex_color.startswith("#"):
            return RGBColor(0, 0, 0)
        hex_color = hex_color.lstrip("#")
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    
    def get_shape_type(shape_type_str: str):
        """文字列からMSO_SHAPEを取得"""
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
            "circle": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "right_arrow": MSO_SHAPE.RIGHT_ARROW,
            "left_arrow": MSO_SHAPE.LEFT_ARROW,
            "up_arrow": MSO_SHAPE.UP_ARROW,
            "down_arrow": MSO_SHAPE.DOWN_ARROW,
        }
        return shape_map.get(shape_type_str, MSO_SHAPE.RECTANGLE)
    
    try:
        # プレゼンテーション作成（16:9）
        prs = Presentation()
        slide_width_inches = 13.333
        slide_height_inches = 7.5
        prs.slide_width = Inches(slide_width_inches)
        prs.slide_height = Inches(slide_height_inches)
        
        # 空のスライドを追加
        blank_layout = prs.slide_layouts[6]  # 空白レイアウト
        slide = prs.slides.add_slide(blank_layout)
        
        # 元画像サイズ（ピクセル）
        original_width = metadata.get("original_size", {}).get("width", 1920)
        original_height = metadata.get("original_size", {}).get("height", 1080)
        
        # ピクセル→インチ変換（96 DPI想定）
        # スライドサイズに合わせてスケーリング
        scale_x = slide_width_inches / (original_width / 96)
        scale_y = slide_height_inches / (original_height / 96)
        scale = min(scale_x, scale_y) * 0.95  # 95%で余白確保
        
        # オフセット計算（中央配置）
        scaled_width = (original_width / 96) * scale
        scaled_height = (original_height / 96) * scale
        offset_x = (slide_width_inches - scaled_width) / 2
        offset_y = (slide_height_inches - scaled_height) / 2
        
        # 背景色設定
        bg_color = metadata.get("background_color")
        if bg_color and bg_color.startswith("#"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(bg_color)
        
        # レイヤーをZ-indexでソートして配置
        layers = sorted(metadata.get("layers", []), key=lambda x: x.get("z_index", 0))
        
        placed_count = 0
        text_count = 0
        shape_count = 0
        image_count = 0
        
        for layer in layers:
            element_type = layer.get("element_type", "image")
            
            # 座標変換（ピクセル→インチ）
            left_inches = offset_x + (layer.get("x", 0) / 96) * scale
            top_inches = offset_y + (layer.get("y", 0) / 96) * scale
            width_inches = (layer.get("width", 100) / 96) * scale
            height_inches = (layer.get("height", 100) / 96) * scale
            
            # 最小サイズを保証（0.1インチ以上）
            width_inches = max(width_inches, 0.1)
            height_inches = max(height_inches, 0.1)
            
            try:
                # テキスト要素 → 編集可能なテキストボックス
                if element_type == "text":
                    content = layer.get("content", "")
                    if content:
                        textbox = slide.shapes.add_textbox(
                            Inches(left_inches),
                            Inches(top_inches),
                            Inches(width_inches),
                            Inches(height_inches)
                        )
                        tf = textbox.text_frame
                        tf.word_wrap = True
                        
                        p = tf.paragraphs[0]
                        p.text = content
                        
                        # フォント設定
                        font_size = layer.get("font_size", 12)
                        font_color = layer.get("font_color", "#000000")
                        font_bold = layer.get("font_bold", False)
                        text_align = layer.get("text_align", "left")
                        
                        run = p.runs[0] if p.runs else p.add_run()
                        run.text = content
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = hex_to_rgb(font_color)
                        run.font.bold = font_bold
                        
                        # テキスト配置
                        align_map = {
                            "left": PP_ALIGN.LEFT,
                            "center": PP_ALIGN.CENTER,
                            "right": PP_ALIGN.RIGHT
                        }
                        p.alignment = align_map.get(text_align, PP_ALIGN.LEFT)
                        
                        text_count += 1
                        placed_count += 1
                        continue
                
                # 図形要素（shape, line, arrow）→ PowerPointネイティブシェイプ
                elif element_type in ["shape", "line", "arrow"]:
                    shape_type_str = layer.get("shape_type", "rectangle")
                    fill_color = layer.get("fill_color")
                    border_color = layer.get("border_color")
                    border_width = layer.get("border_width", 1)
                    
                    # 線の場合
                    if element_type == "line" or shape_type_str == "line":
                        # 線はコネクタとして追加（簡易版：四角形で細い線を表現）
                        line_shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(left_inches),
                            Inches(top_inches),
                            Inches(width_inches),
                            Inches(max(height_inches, 0.02))  # 最小2ピクセル相当
                        )
                        line_shape.fill.solid()
                        if border_color:
                            line_shape.fill.fore_color.rgb = hex_to_rgb(border_color)
                        else:
                            line_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
                        line_shape.line.fill.background()  # 線なし
                        shape_count += 1
                        placed_count += 1
                        continue
                    
                    # 矢印の場合
                    if element_type == "arrow" or shape_type_str == "arrow":
                        shape_type = MSO_SHAPE.RIGHT_ARROW
                    else:
                        shape_type = get_shape_type(shape_type_str)
                    
                    shape = slide.shapes.add_shape(
                        shape_type,
                        Inches(left_inches),
                        Inches(top_inches),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # 塗りつぶし設定
                    if fill_color:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
                    else:
                        shape.fill.background()  # 透明
                    
                    # 線（境界）設定
                    if border_color:
                        shape.line.color.rgb = hex_to_rgb(border_color)
                        shape.line.width = Pt(border_width or 1)
                    else:
                        shape.line.fill.background()  # 線なし
                    
                    shape_count += 1
                    placed_count += 1
                    continue
                
                # 画像/アイコン/その他 → PNG画像として配置
                image_path = layer.get("image_path")
                if image_path and Path(image_path).exists():
                    slide.shapes.add_picture(
                        image_path,
                        Inches(left_inches),
                        Inches(top_inches),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    image_count += 1
                    placed_count += 1
                
            except Exception as e:
                logger.warning(f"レイヤー配置エラー: {layer.get('layer_id')} - {e}")
                # フォールバック: 画像として配置を試みる
                image_path = layer.get("image_path")
                if image_path and Path(image_path).exists():
                    try:
                        slide.shapes.add_picture(
                            image_path,
                            Inches(left_inches),
                            Inches(top_inches),
                            Inches(width_inches),
                            Inches(height_inches)
                        )
                        image_count += 1
                        placed_count += 1
                    except:
                        pass
                continue
        
        # PPTX保存
        output_filename = f"{file_id}_layers.pptx"
        output_path = OUTPUT_DIR / output_filename
        prs.save(str(output_path))
        
        return {
            "file_id": file_id,
            "success": True,
            "pptx_path": str(output_path),
            "download_url": f"/api/layers/{file_id}/download",
            "layer_count": placed_count,
            "text_count": text_count,
            "shape_count": shape_count,
            "image_count": image_count,
            "message": f"{placed_count}個のレイヤーをPPTXに配置しました（テキスト: {text_count}, 図形: {shape_count}, 画像: {image_count}）"
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPTX生成エラー: {str(e)}")


@app.get("/api/layers/{file_id}/download")
async def download_layers_pptx(file_id: str):
    """分解レイヤーから生成したPPTXをダウンロード"""
    output_path = OUTPUT_DIR / f"{file_id}_layers.pptx"
    
    if not output_path.exists():
        raise HTTPException(status_code=404, detail="PPTXファイルが見つかりません")
    
    return FileResponse(
        path=str(output_path),
        filename=f"layers_{file_id}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
